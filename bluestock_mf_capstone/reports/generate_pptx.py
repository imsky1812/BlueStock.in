"""
generate_pptx.py
================
BlueStock Fintech Internship — Capstone Project I
Generates the final Presentation.pptx (D7)

Run:
  pip install python-pptx
  python reports/generate_pptx.py

Output: reports/Presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ─── Paths ────────────────────────────────────────────────────────────────────
OUT_DIR = Path(__file__).resolve().parent
OUT_PATH = OUT_DIR / "Presentation.pptx"

# ─── Color Palette ────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1e, 0x3a, 0x8a)
TEAL       = RGBColor(0x0d, 0x94, 0x88)
AMBER      = RGBColor(0xf5, 0x9e, 0x0b)
RED        = RGBColor(0xdc, 0x26, 0x26)
GREEN      = RGBColor(0x16, 0xa3, 0x4a)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_GREY = RGBColor(0xf1, 0xf5, 0xf9)
DARK       = RGBColor(0x0f, 0x17, 0x2a)
MID        = RGBColor(0x47, 0x55, 0x69)


def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ─── Helpers ──────────────────────────────────────────────────────────────────
def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                 wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─── Slide Definitions ────────────────────────────────────────────────────────

SLIDES = [
    # slide_func, title
    "slide_cover",
    "slide_agenda",
    "slide_project_overview",
    "slide_data_sources",
    "slide_etl",
    "slide_eda",
    "slide_performance",
    "slide_risk",
    "slide_var_cvar",
    "slide_cohort",
    "slide_sip_continuity",
    "slide_monte_carlo",
    "slide_efficient_frontier",
    "slide_recommendations",
    "slide_thankyou",
]


def build_cover(slide, prs):
    set_bg(slide, DARK_BLUE)
    # Top accent bar
    add_rect(slide, 0, 0, 10, 0.12, TEAL)
    # Bottom accent bar
    add_rect(slide, 0, 7.38, 10, 0.12, AMBER)
    # Main title
    add_text_box(slide, "Mutual Fund Industry Analytics", 0.8, 1.2, 8.4, 1.2,
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Capstone Project I — BlueStock Fintech Internship", 0.8, 2.5, 8.4, 0.6,
                 font_size=20, bold=False, color=RGBColor(0xa5, 0xb4, 0xfc), align=PP_ALIGN.CENTER)
    add_rect(slide, 2.5, 3.3, 5.0, 0.05, TEAL)
    add_text_box(slide, "👥 Sarvesh Kumar Yadav", 0.8, 3.6, 8.4, 0.5,
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "📅 Submission: July 10, 2026   |   📊 Data: Jan 2020 – May 2025", 0.8, 4.2, 8.4, 0.5,
                 font_size=13, color=RGBColor(0xcb, 0xd5, 0xe1), align=PP_ALIGN.CENTER)
    add_text_box(slide, "BlueStock Fintech · Batch 2026", 0.8, 6.6, 8.4, 0.4,
                 font_size=11, color=RGBColor(0x64, 0x74, 0x8b), align=PP_ALIGN.CENTER)


def build_agenda(slide, prs):
    set_bg(slide, LIGHT_GREY)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Agenda", 0.4, 0.2, 9.2, 0.7, font_size=28, bold=True, color=WHITE)
    items = [
        ("01", "Project Overview & Scope"),
        ("02", "Data Architecture & ETL Pipeline (D1, D2)"),
        ("03", "Exploratory Data Analysis (D3)"),
        ("04", "Fund Performance Analytics (D4)"),
        ("05", "VaR, CVaR & Risk Metrics (D6)"),
        ("06", "Investor Cohort & SIP Analysis"),
        ("07", "Monte Carlo Simulation 5-Year Projection [B3]"),
        ("08", "Markowitz Efficient Frontier [B4]"),
        ("09", "Dashboard & Recommendations (D5, D7)"),
    ]
    for i, (num, text) in enumerate(items):
        row = i // 2
        col = i % 2
        x = 0.3 + col * 4.9
        y = 1.3 + row * 1.0
        add_rect(slide, x, y, 0.45, 0.45, DARK_BLUE)
        add_text_box(slide, num, x, y, 0.45, 0.45, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, text, x + 0.55, y + 0.02, 4.0, 0.4, font_size=13, color=DARK)


def build_project_overview(slide, prs):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Project Overview & Scope", 0.4, 0.2, 9.2, 0.7, font_size=26, bold=True, color=WHITE)
    kpis = [("40", "Fund Schemes"), ("10", "AMCs"), ("₹67.2L Cr", "Industry AUM"), ("5 Years", "NAV Data")]
    for i, (val, lbl) in enumerate(kpis):
        x = 0.3 + i * 2.35
        add_rect(slide, x, 1.3, 2.1, 1.3, DARK_BLUE)
        add_text_box(slide, val, x, 1.4, 2.1, 0.7, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x, 2.1, 2.1, 0.4, font_size=11, color=RGBColor(0xa5, 0xb4, 0xfc), align=PP_ALIGN.CENTER)
    bullets = [
        "✅  D1 ETL pipeline · D2 SQLite DB · D3 EDA notebook",
        "✅  D4 Performance metrics · D5 Dashboard · D6 Advanced analytics",
        "✅  D7 Final report + slides · B2 Streamlit · B3 Monte Carlo · B4 Frontier",
        "📦  Data: 10 CSV datasets →  cleaned → SQLite star schema",
        "🛠️  Stack: Python 3.11 · Pandas · SQLAlchemy · Streamlit · Plotly · SciPy",
    ]
    for i, b in enumerate(bullets):
        add_text_box(slide, b, 0.4, 2.8 + i * 0.65, 9.2, 0.55, font_size=13, color=DARK)


def build_data_sources(slide, prs):
    set_bg(slide, LIGHT_GREY)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Data Sources — 10 Datasets", 0.4, 0.2, 9.2, 0.7, font_size=26, bold=True, color=WHITE)
    rows = [
        ("01", "Fund Master", "40 rows", "Scheme metadata, risk grade, expense ratio"),
        ("02", "NAV History", "~50,000 rows", "Daily NAV per fund (5-year history)"),
        ("03", "AUM by Fund House", "180 rows", "Monthly AUM breakdown across 10 AMCs"),
        ("04", "Monthly SIP Inflows", "60 rows", "Industry SIP statistics 2020–2025"),
        ("05", "Category Inflows", "~300 rows", "Net inflows by fund category per month"),
        ("06", "Folio Count", "60 rows", "Total industry folios (equity/debt/hybrid)"),
        ("07", "Scheme Performance", "40 rows", "Pre-computed Sharpe, Alpha, Beta"),
        ("08", "Investor Transactions", "~100,000 rows", "SIP/Lumpsum/Redemption records"),
        ("09", "Portfolio Holdings", "~500 rows", "Stock-level holdings with sector weights"),
        ("10", "Benchmark Indices", "~6,000 rows", "NIFTY 50, NIFTY 500 daily OHLCV"),
    ]
    for i, (num, name, size, desc) in enumerate(rows):
        row_y = 1.2 + i * 0.58
        col = 0 if i < 5 else 1
        row_i = i if i < 5 else i - 5
        x = 0.2 + col * 4.9
        y = 1.2 + row_i * 1.15
        add_rect(slide, x, y, 0.4, 0.4, TEAL)
        add_text_box(slide, num, x, y+0.04, 0.4, 0.32, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"{name} ({size})", x+0.45, y+0.02, 4.2, 0.25, font_size=12, bold=True, color=DARK_BLUE)
        add_text_box(slide, desc, x+0.45, y+0.28, 4.2, 0.25, font_size=10, color=MID)


def build_etl(slide, prs):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "ETL Pipeline Architecture (D1 + D2)", 0.4, 0.2, 9.2, 0.7, font_size=26, bold=True, color=WHITE)
    stages = [("📥 EXTRACT", "Load 10 raw CSVs\nValidate file existence\npathlib.Path — no hard-coded paths"),
              ("🔧 TRANSFORM", "Date parsing & type coercion\nDeduplication\nNAV ffill() for weekends/holidays"),
              ("📤 LOAD", "Build dim_date table\nLoad 11 SQLite tables\nSQLAlchemy engine"),
              ("📋 REPORT", "Null/dup/AMFI stats\nData quality report\nETL log file")]
    colors = [DARK_BLUE, TEAL, RGBColor(0x8b, 0x5c, 0xf6), AMBER]
    for i, (title, body) in enumerate(stages):
        x = 0.3 + i * 2.35
        add_rect(slide, x, 1.2, 2.1, 2.8, colors[i])
        add_text_box(slide, title, x, 1.3, 2.1, 0.5, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body, x+0.1, 1.9, 1.9, 2.0, font_size=10, color=WHITE)
        if i < 3:
            add_text_box(slide, "→", x+2.05, 2.4, 0.3, 0.4, font_size=22, bold=True, color=DARK_BLUE)
    add_text_box(slide, "Key Design Principles:", 0.4, 4.2, 9, 0.3, font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, "✅ pathlib.Path (no absolute paths)   ✅ ffill() after reindex to business days   ✅ CAGR uses 252 trading days   ✅ Full try/except error handling",
                 0.4, 4.55, 9.2, 0.5, font_size=12, color=DARK)


def build_eda(slide, prs):
    set_bg(slide, LIGHT_GREY)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Exploratory Data Analysis — Key Findings (D3)", 0.4, 0.2, 9.2, 0.7, font_size=24, bold=True, color=WHITE)
    findings = [
        ("📈", "SIP Tripled in 5 Years", "Monthly SIP inflows grew from ₹7,300 Cr (FY21) to ₹21,400 Cr (FY25) — a 3× increase driven by digital adoption."),
        ("🏙️", "Tier-1 City Dominance", "68% of invested capital comes from Mumbai, Delhi, Bengaluru, Hyderabad, and Chennai — Tier-2/3 is an untapped market."),
        ("⚖️", "Equity Leads AUM", "Equity funds hold 58% of AUM, up from 48% in 2020, reflecting bullish market sentiment and growing retail participation."),
        ("🔢", "Fund Universe", "40 schemes across 10 AMCs. Axis, HDFC, SBI, ICICI, and Kotak collectively hold 78% of total AUM in the universe."),
    ]
    for i, (icon, title, body) in enumerate(findings):
        col, row = i % 2, i // 2
        x = 0.3 + col * 4.85
        y = 1.3 + row * 2.8
        add_rect(slide, x, y, 4.55, 2.5, WHITE)
        add_text_box(slide, icon, x+0.15, y+0.2, 0.6, 0.6, font_size=24, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x+0.8, y+0.25, 3.6, 0.4, font_size=13, bold=True, color=DARK_BLUE)
        add_text_box(slide, body, x+0.15, y+0.8, 4.2, 1.5, font_size=11, color=MID)


def build_performance(slide, prs):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Fund Performance Analytics (D4)", 0.4, 0.2, 9.2, 0.7, font_size=26, bold=True, color=WHITE)
    add_text_box(slide, "Top 5 Funds by Sharpe Ratio", 0.4, 1.2, 4.5, 0.4, font_size=14, bold=True, color=DARK_BLUE)
    headers = ["Fund", "CAGR", "Sharpe", "Beta"]
    col_x = [0.4, 4.0, 5.3, 6.5]
    col_w = [3.4, 1.0, 1.0, 1.0]
    add_rect(slide, 0.3, 1.65, 7.2, 0.4, DARK_BLUE)
    for j, h in enumerate(headers):
        add_text_box(slide, h, col_x[j], 1.7, col_w[j], 0.32, font_size=11, bold=True, color=WHITE)
    top5 = [
        ("ABSL Liquid Fund", "5.2%", "3.21", "0.04"),
        ("ICICI Pru Liquid Fund", "5.1%", "3.18", "0.03"),
        ("SBI Small Cap – Direct", "28.4%", "1.52", "0.82"),
        ("HDFC Top 100 – Direct", "21.2%", "1.42", "0.93"),
        ("Axis Midcap – Direct", "24.7%", "1.38", "0.79"),
    ]
    row_colors = [LIGHT_GREY, WHITE, LIGHT_GREY, WHITE, LIGHT_GREY]
    for i, (fund, cagr, sharpe, beta) in enumerate(top5):
        y = 2.1 + i * 0.65
        add_rect(slide, 0.3, y, 7.2, 0.6, row_colors[i])
        for j, val in enumerate([fund, cagr, sharpe, beta]):
            add_text_box(slide, val, col_x[j], y+0.12, col_w[j], 0.38, font_size=11, color=DARK)
    add_text_box(slide, "Key Formulas Used:", 7.8, 1.2, 2.0, 0.3, font_size=12, bold=True, color=DARK_BLUE)
    formulas = ["CAGR = (NAV_end/NAV_start)^(252/days)−1", "Sharpe = (r−rf)/σ × √252", "Beta = Cov(r,rm)/Var(rm)", "Alpha = r−[rf+β(rm−rf)]", "rf = 6.5%/252 daily"]
    for i, f in enumerate(formulas):
        add_rect(slide, 7.7, 1.6 + i * 0.95, 2.1, 0.8, LIGHT_GREY)
        add_text_box(slide, f, 7.75, 1.65 + i * 0.95, 2.0, 0.7, font_size=9, color=DARK_BLUE)
    add_text_box(slide, "⚠️ Liquid funds rank highest by Sharpe due to near-zero volatility — not highest absolute return.", 0.4, 5.3, 9.2, 0.5, font_size=11, color=RED)


def build_risk(slide, prs):
    set_bg(slide, LIGHT_GREY)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Sector Concentration — HHI Analysis (D6)", 0.4, 0.2, 9.2, 0.7, font_size=26, bold=True, color=WHITE)
    add_text_box(slide, "HHI = Σ(sector_weight)²   |   HHI > 2500 = Concentrated   |   HHI < 1500 = Diversified",
                 0.4, 1.2, 9.2, 0.4, font_size=12, color=DARK_BLUE)
    bars = [("Axis Bluechip Fund", 2968, 3000, RED),
            ("SBI Bluechip Fund", 2541, 3000, RED),
            ("HDFC Top 100 Fund", 2312, 3000, AMBER),
            ("Axis Midcap Fund", 1820, 3000, AMBER),
            ("UTI Mid Cap Fund", 1240, 3000, GREEN)]
    for i, (name, hhi, maxv, color) in enumerate(bars):
        y = 1.8 + i * 0.9
        add_text_box(slide, name, 0.4, y, 3.0, 0.4, font_size=12, color=DARK)
        bar_w = 5.0 * hhi / maxv
        add_rect(slide, 3.6, y+0.05, bar_w, 0.35, color)
        add_text_box(slide, str(hhi), 3.6 + bar_w + 0.1, y+0.05, 0.8, 0.35, font_size=11, bold=True, color=DARK)
    add_text_box(slide, "💡 Large-cap 'blue chip' funds are paradoxically MORE concentrated (heavy BFSI/IT bets) than mid/small cap diversified funds.",
                 0.4, 6.4, 9.2, 0.6, font_size=12, color=DARK_BLUE)


def build_var_cvar(slide, prs):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Historical VaR 95% & CVaR (Expected Shortfall)", 0.4, 0.2, 9.2, 0.7, font_size=24, bold=True, color=WHITE)
    add_text_box(slide, "VaR 95% = 5th percentile of daily returns  |  CVaR = Mean of returns ≤ VaR",
                 0.4, 1.2, 9.2, 0.4, font_size=12, color=MID)
    data = [("Small Cap", "SBI Small Cap Fund", "–2.69%", "–3.82%", RED),
            ("Mid Cap",   "Axis Midcap Fund",   "–2.21%", "–3.14%", AMBER),
            ("Large Cap", "HDFC Top 100",        "–1.74%", "–2.48%", AMBER),
            ("Hybrid",    "ICICI Pru Balanced",  "–0.92%", "–1.31%", GREEN),
            ("Liquid",    "ABSL Liquid Fund",    "–0.02%", "–0.03%", GREEN)]
    add_rect(slide, 0.3, 1.7, 9.4, 0.42, DARK_BLUE)
    for j, h in enumerate(["Category", "Fund", "95% VaR", "95% CVaR"]):
        add_text_box(slide, h, [0.35, 1.8, 5.5, 7.3][j], 1.75, [1.3, 3.5, 1.5, 1.5][j], 0.32, font_size=11, bold=True, color=WHITE)
    for i, (cat, fund, var, cvar, color) in enumerate(data):
        y = 2.18 + i * 0.78
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        add_rect(slide, 0.3, y, 9.4, 0.72, bg)
        add_rect(slide, 0.3, y, 0.08, 0.72, color)
        add_text_box(slide, cat,  0.45, y+0.17, 1.3,  0.38, font_size=11, bold=True, color=color)
        add_text_box(slide, fund, 1.8,  y+0.17, 3.5,  0.38, font_size=11, color=DARK)
        add_text_box(slide, var,  5.5,  y+0.17, 1.5,  0.38, font_size=12, bold=True, color=RED)
        add_text_box(slide, cvar, 7.3,  y+0.17, 1.5,  0.38, font_size=12, bold=True, color=RED)


def build_cohort(slide, prs):
    set_bg(slide, LIGHT_GREY)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Investor Cohort Analysis (D6)", 0.4, 0.2, 9.2, 0.7, font_size=26, bold=True, color=WHITE)
    cohorts = [("2020", 312, 8240, "28.4 Cr"), ("2021", 891, 9100, "86.2 Cr"),
               ("2022", 1847, 10200, "198.7 Cr"), ("2023", 2760, 10760, "312.4 Cr"),
               ("2024", 4803, 10997, "225.8 Cr"), ("2025*", 197, 13505, "19.0 Cr")]
    max_inv = max(c[1] for c in cohorts)
    for i, (yr, inv, sip, total) in enumerate(cohorts):
        x = 0.4 + i * 1.55
        bar_h = 3.5 * inv / max_inv
        y_base = 5.0
        add_rect(slide, x, y_base - bar_h, 1.2, bar_h, DARK_BLUE if yr != "2024" else TEAL)
        add_text_box(slide, str(inv), x, y_base - bar_h - 0.35, 1.2, 0.3, font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide, yr,       x, y_base + 0.05,          1.2, 0.3, font_size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"Avg SIP\n₹{sip:,}", x, y_base + 0.4, 1.2, 0.55, font_size=9, color=MID, align=PP_ALIGN.CENTER)
    add_text_box(slide, "New Investors per Cohort Year (bar height ∝ count)  |  *2025 data through May only",
                 0.4, 5.55, 9.2, 0.35, font_size=10, color=MID)
    add_text_box(slide, "📌 2025 cohort: smallest count, highest avg SIP (₹13,505) → wealthier/more sophisticated new investors",
                 0.4, 6.6, 9.2, 0.5, font_size=12, bold=True, color=DARK_BLUE)


def build_sip_continuity(slide, prs):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "SIP Continuity & Churn Risk Analysis (D6)", 0.4, 0.2, 9.2, 0.7, font_size=24, bold=True, color=WHITE)
    add_text_box(slide, "Investors with ≥ 6 SIP transactions", 0.4, 1.2, 9, 0.4, font_size=14, bold=True, color=DARK_BLUE)
    metrics = [("1,362", "Total Investors\n(6+ SIPs)"), ("30", "Active\n(gap ≤ 35 days)"),
               ("1,332", "At-Risk\n(gap > 35 days)"), ("97.8%", "Churn Risk\nRate")]
    colors2 = [DARK_BLUE, GREEN, RED, RED]
    for i, (val, lbl) in enumerate(metrics):
        x = 0.4 + i * 2.35
        add_rect(slide, x, 1.7, 2.1, 1.5, colors2[i])
        add_text_box(slide, val, x, 1.8, 2.1, 0.7, font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x, 2.55, 2.1, 0.6, font_size=10, color=RGBColor(0xa5, 0xb4, 0xfc) if colors2[i]==DARK_BLUE else WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Definition: avg_gap_days = mean(diff between consecutive SIP dates) for each investor.\nStatus 'at-risk' if avg_gap > 35 days (standard monthly SIP = ~30 days).",
                 0.4, 3.4, 9.2, 0.8, font_size=12, color=MID)
    add_text_box(slide, "🚨 Critical: 97.8% churn risk rate signals systematic SIP mandate/payment execution failures.\n   Even a 20% recovery = ₹50+ Cr additional annual SIP inflows.",
                 0.4, 4.4, 9.2, 0.8, font_size=13, bold=True, color=RED)
    add_text_box(slide, "💡 Recommended Action: T+3 day SMS/push notification for missed SIP + automated mandate retry.",
                 0.4, 5.5, 9.2, 0.6, font_size=13, color=DARK_BLUE)


def build_monte_carlo(slide, prs):
    set_bg(slide, LIGHT_GREY)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "[B3] Monte Carlo Simulation — 5-Year NAV Projection", 0.4, 0.2, 9.2, 0.7, font_size=24, bold=True, color=WHITE)
    add_text_box(slide, "Method: Geometric Brownian Motion (GBM)   |   1,000 simulation paths   |   1,260 trading days (5 years)",
                 0.4, 1.2, 9.2, 0.35, font_size=11, color=MID)
    add_rect(slide, 0.3, 1.65, 9.4, 0.4, DARK_BLUE)
    headers = ["Fund", "Start NAV", "Bear (P5)", "Median (P50)", "Bull (P95)", "Expected CAGR"]
    xs = [0.35, 3.1, 4.4, 5.7, 7.0, 8.3]
    ws = [2.6, 1.0, 1.0, 1.0, 1.0, 1.0]
    for j, h in enumerate(headers):
        add_text_box(slide, h, xs[j], 1.7, ws[j], 0.32, font_size=10, bold=True, color=WHITE)
    data = [("SBI Small Cap – Direct",   "₹158.4", "₹89.2",   "₹312.8",   "₹618.4",  "14.6%"),
            ("Axis Midcap – Direct",     "₹98.7",  "₹58.3",   "₹187.4",   "₹356.8",  "13.7%"),
            ("HDFC Top 100 – Direct",    "₹892.6", "₹562.4",  "₹1,624.8", "₹2,958.3","12.7%"),
            ("SBI Bluechip – Direct",    "₹74.3",  "₹46.8",   "₹133.2",   "₹238.9",  "12.4%"),
            ("ICICI Pru Bluechip",       "₹102.8", "₹64.7",   "₹183.6",   "₹327.4",  "12.3%")]
    for i, row in enumerate(data):
        y = 2.12 + i * 0.78
        bg = WHITE if i % 2 == 0 else LIGHT_GREY
        add_rect(slide, 0.3, y, 9.4, 0.72, bg)
        for j, val in enumerate(row):
            add_text_box(slide, val, xs[j], y+0.18, ws[j], 0.38, font_size=10, color=DARK)
    add_text_box(slide, "⚠️ P5–P95 spread for equity funds is 3–5× → long holding period (≥5 years) essential to absorb volatility.",
                 0.4, 6.2, 9.2, 0.5, font_size=12, color=RED)


def build_efficient_frontier(slide, prs):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "[B4] Markowitz Efficient Frontier — Portfolio Optimisation", 0.4, 0.2, 9.2, 0.7, font_size=22, bold=True, color=WHITE)
    add_text_box(slide, "5 equity funds · 5,000 random portfolios · SciPy SLSQP optimiser · Risk-free rate: 6.5%",
                 0.4, 1.2, 9.2, 0.35, font_size=11, color=MID)
    for i, (label, ret, vol, sharpe, weights, color) in enumerate([
        ("⭐ Max Sharpe Portfolio",    "21.4%", "13.8%", "1.11",
         "Axis Midcap 31% · HDFC Top100 29% · SBI Small Cap 19%", TEAL),
        ("🛡️ Min Variance Portfolio", "17.8%", "11.2%", "1.01",
         "HDFC Top100 42% · SBI Bluechip 32% · ICICI Bluechip 13%", DARK_BLUE),
        ("⚖️ Equal-Weight (Naive)",   "19.3%", "14.9%", "0.86",
         "20% each fund", MID),
    ]):
        y = 1.65 + i * 1.85
        add_rect(slide, 0.3, y, 9.4, 1.65, LIGHT_GREY)
        add_rect(slide, 0.3, y, 0.08, 1.65, color)
        add_text_box(slide, label, 0.5, y+0.15, 4.5, 0.4, font_size=13, bold=True, color=color)
        stats = f"Return: {ret}   |   Volatility: {vol}   |   Sharpe: {sharpe}"
        add_text_box(slide, stats, 0.5, y+0.55, 9.0, 0.35, font_size=12, color=DARK)
        add_text_box(slide, f"Weights: {weights}", 0.5, y+0.95, 9.0, 0.35, font_size=11, color=MID)
    add_text_box(slide, "✅ Max Sharpe Portfolio achieves 29% higher Sharpe than equal-weight allocation.",
                 0.4, 7.2, 9.2, 0.4, font_size=13, bold=True, color=GREEN)


def build_recommendations(slide, prs):
    set_bg(slide, LIGHT_GREY)
    add_rect(slide, 0, 0, 10, 1.1, DARK_BLUE)
    add_text_box(slide, "Strategic Recommendations", 0.4, 0.2, 9.2, 0.7, font_size=26, bold=True, color=WHITE)
    recs = [
        ("01", TEAL,    "MPT-Guided Portfolio",     "Replace equal-weight defaults with Max Sharpe allocation. +29% Sharpe improvement."),
        ("02", RED,     "SIP Nudge Programme",      "T+3 SMS/push for missed SIPs. 97.8% at-risk rate = ₹50+ Cr recoverable annual AUM."),
        ("03", AMBER,   "Tier-2/3 Expansion",       "Only 10% of capital from Tier-3. Lower min SIP + vernacular campaigns = new TAM."),
        ("04", DARK_BLUE,"Risk Suitability Check",  "First-time 2024 investors over-allocated to Small Cap (VaR –2.69%). Add onboarding guardrails."),
        ("05", GREEN,   "Monte Carlo Planning Tool","Integrate 5-year P5/P50/P95 projections into investor app to reduce premature redemptions."),
    ]
    for i, (num, color, title, body) in enumerate(recs):
        col, row = i % 2, i // 2
        if i == 4: col, row = 0, 2
        x = 0.3 + col * 4.85
        y = 1.3 + row * 2.0
        add_rect(slide, x, y, 4.55, 1.8, WHITE)
        add_rect(slide, x, y, 0.07, 1.8, color)
        add_rect(slide, x+0.15, y+0.15, 0.45, 0.45, color)
        add_text_box(slide, num, x+0.15, y+0.17, 0.45, 0.38, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x+0.75, y+0.2, 3.6, 0.4, font_size=13, bold=True, color=DARK_BLUE)
        add_text_box(slide, body, x+0.75, y+0.65, 3.6, 1.0, font_size=10, color=MID)


def build_thankyou(slide, prs):
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.12, AMBER)
    add_rect(slide, 0, 7.38, 10, 0.12, TEAL)
    add_text_box(slide, "Thank You", 0, 1.5, 10, 1.5, font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Questions & Discussion", 0, 3.2, 10, 0.6, font_size=22, color=RGBColor(0xa5, 0xb4, 0xfc), align=PP_ALIGN.CENTER)
    add_rect(slide, 3, 4.0, 4, 0.05, TEAL)
    add_text_box(slide, "Sarvesh Kumar Yadav", 0, 4.3, 10, 0.5, font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "BlueStock Fintech Internship · Batch 2026 · July 10, 2026", 0, 4.85, 10, 0.4, font_size=13, color=RGBColor(0x64, 0x74, 0x8b), align=PP_ALIGN.CENTER)
    deliverables = "D1 ETL  ·  D2 SQLite  ·  D3 EDA  ·  D4 Performance  ·  D5 Dashboard  ·  D6 Analytics  ·  D7 Report  ·  B2 Streamlit  ·  B3 Monte Carlo  ·  B4 Frontier"
    add_text_box(slide, deliverables, 0.4, 6.0, 9.2, 0.5, font_size=10, color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.CENTER)


SLIDE_BUILDERS = [
    build_cover, build_agenda, build_project_overview, build_data_sources,
    build_etl, build_eda, build_performance, build_risk, build_var_cvar,
    build_cohort, build_sip_continuity, build_monte_carlo, build_efficient_frontier,
    build_recommendations, build_thankyou,
]


def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for i, builder in enumerate(SLIDE_BUILDERS, 1):
        slide = prs.slides.add_slide(blank_layout)
        builder(slide, prs)
        print(f"  [OK] Slide {i:02d}/{len(SLIDE_BUILDERS)} - {builder.__name__}")

    prs.save(str(OUT_PATH))
    print(f"\n[DONE] Presentation saved: {OUT_PATH}")
    print(f"       {len(SLIDE_BUILDERS)} slides total")


if __name__ == "__main__":
    main()
